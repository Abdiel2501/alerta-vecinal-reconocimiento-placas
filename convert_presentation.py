import sys
import os
import re
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()
# Set 16:9 Widescreen dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color constants
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_GRAY_DARK = RGBColor(60, 60, 60)
COLOR_GRAY_LIGHT = RGBColor(248, 250, 252)
COLOR_CARD_BG = RGBColor(255, 255, 255)
COLOR_TECH_BG = RGBColor(248, 250, 252)
COLOR_BORDER = RGBColor(0, 0, 0)

FONT_TITLE = "Arial Black"
FONT_BODY = "Arial"
FONT_CODE = "Consolas"

# Load and parse HTML
html_path = "presentacion_proyecto.html"
with open(html_path, "r", encoding="utf-8") as f:
    soup = BeautifulSoup(f.read(), "html.parser")

slides_html = soup.find_all("section", class_="slide")

def add_neobrutalist_card(slide, left, top, width, height, fill_color=COLOR_CARD_BG, border_color=COLOR_BORDER):
    # Shadow (black rectangle slightly offset)
    shadow_offset = Inches(0.12)
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        left + shadow_offset, 
        top + shadow_offset, 
        width, 
        height
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = COLOR_BLACK
    shadow.line.fill.background() # No border for shadow
    
    # Foreground Card
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        left, 
        top, 
        width, 
        height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(3)
    
    return card

def clean_text(text):
    if not text:
        return ""
    # Strip whitespace, replace multiple spaces with single space
    return re.sub(r'\s+', ' ', text).strip()

def format_text_frame(tf):
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.25)
    tf.margin_bottom = Inches(0.25)

# Process slides
for index, slide_node in enumerate(slides_html):
    # Determine slide layout (blank layout is 6 in python-pptx)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Check if Cover slide (Slide 1 has cover-container)
    cover_container = slide_node.find("div", class_="cover-container")
    if cover_container:
        # Cover Slide Layout
        # Draw a nice clean cover
        # Background is white, add some bold elements
        
        # Draw top accent bar
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLOR_BLACK
        top_bar.line.fill.background()
        
        # Subtitle/Tag
        tag_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(0.5))
        tf = tag_box.text_frame
        format_text_frame(tf)
        p = tf.paragraphs[0]
        p.text = "🛡️ EXPLICACIÓN TÉCNICA DE CÓDIGO"
        p.font.name = FONT_TITLE
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLACK
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
        tf = title_box.text_frame
        format_text_frame(tf)
        p = tf.paragraphs[0]
        p.text = "SISTEMA DE DETECCIÓN\nY RECONOCIMIENTO"
        p.font.name = FONT_TITLE
        p.font.size = Pt(46)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLACK
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.333), Inches(1.0))
        tf = sub_box.text_frame
        format_text_frame(tf)
        p = tf.paragraphs[0]
        p.text = "¿Cómo funciona nuestro código, qué librerías utilizamos y por qué tomamos cada decisión de ingeniería?"
        p.font.name = FONT_BODY
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_GRAY_DARK
        
        # Meta info cards
        # Servidor info
        add_neobrutalist_card(slide, Inches(1.0), Inches(5.8), Inches(3.0), Inches(1.0))
        meta1_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(3.0), Inches(1.0))
        tf1 = meta1_box.text_frame
        format_text_frame(tf1)
        p1_label = tf1.paragraphs[0]
        p1_label.text = "SERVIDOR"
        p1_label.font.name = FONT_TITLE
        p1_label.font.size = Pt(11)
        p1_label.font.color.rgb = COLOR_GRAY_DARK
        p1_val = tf1.add_paragraph()
        p1_val.text = "Python Backend"
        p1_val.font.name = FONT_BODY
        p1_val.font.size = Pt(16)
        p1_val.font.bold = True
        p1_val.font.color.rgb = COLOR_BLACK
        
        # Cliente info
        add_neobrutalist_card(slide, Inches(4.5), Inches(5.8), Inches(3.0), Inches(1.0))
        meta2_box = slide.shapes.add_textbox(Inches(4.5), Inches(5.8), Inches(3.0), Inches(1.0))
        tf2 = meta2_box.text_frame
        format_text_frame(tf2)
        p2_label = tf2.paragraphs[0]
        p2_label.text = "CLIENTE"
        p2_label.font.name = FONT_TITLE
        p2_label.font.size = Pt(11)
        p2_label.font.color.rgb = COLOR_GRAY_DARK
        p2_val = tf2.add_paragraph()
        p2_val.text = "Flutter UI"
        p2_val.font.name = FONT_BODY
        p2_val.font.size = Pt(16)
        p2_val.font.bold = True
        p2_val.font.color.rgb = COLOR_BLACK
        
    else:
        # Standard Slide layout
        # Top Progress Bar
        progress_width = (index + 1) / len(slides_html) * 13.333
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(progress_width), Inches(0.12))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLOR_BLACK
        top_bar.line.fill.background()
        
        # Title of the slide
        title_h2 = slide_node.find("h2", class_="slide-title")
        title_text = ""
        if title_h2:
            span = title_h2.find("span")
            title_text = span.get_text() if span else title_h2.get_text()
        else:
            title_text = slide_node.get("data-title", "Diapositiva")
            
        title_text = clean_text(title_text).upper()
        
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
        tf = title_box.text_frame
        format_text_frame(tf)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_TITLE
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLACK
        
        # Left and Right Columns
        cards = slide_node.find_all("div", class_="card-bordered")
        
        col_width = Inches(5.8)
        col_height = Inches(5.0)
        top_y = Inches(1.3)
        left_x = Inches(0.6)
        right_x = Inches(6.9)
        
        # Slide content blocks
        for card_idx, card_node in enumerate(cards[:2]):
            x_pos = left_x if card_idx == 0 else right_x
            
            # Create neo-brutalism card
            add_neobrutalist_card(slide, x_pos, top_y, col_width, col_height)
            
            # Card contents container
            content_box = slide.shapes.add_textbox(x_pos, top_y, col_width, col_height)
            tf_card = content_box.text_frame
            format_text_frame(tf_card)
            
            # Parse components inside card
            h3_node = card_node.find("h3")
            if h3_node:
                p_h3 = tf_card.paragraphs[0]
                p_h3.text = clean_text(h3_node.get_text()).upper()
                p_h3.font.name = FONT_TITLE
                p_h3.font.size = Pt(16)
                p_h3.font.color.rgb = COLOR_BLACK
                # Add spacing below header
                p_h3.space_after = Pt(10)
            
            # Check for Paragraphs / Lists
            ul_node = card_node.find("ul")
            if ul_node:
                for li in ul_node.find_all("li"):
                    p_li = tf_card.add_paragraph()
                    # Keep lists as bullet points
                    p_li.level = 0
                    p_li.space_after = Pt(6)
                    
                    # See if there is a strong element
                    strong = li.find("strong")
                    if strong:
                        # Add bold part
                        run_bold = p_li.add_run()
                        run_bold.text = "■ " + clean_text(strong.get_text()) + " "
                        run_bold.font.bold = True
                        run_bold.font.name = FONT_BODY
                        run_bold.font.size = Pt(13)
                        run_bold.font.color.rgb = COLOR_BLACK
                        
                        # Add rest
                        rest_text = li.get_text().replace(strong.get_text(), "", 1).strip()
                        run_rest = p_li.add_run()
                        run_rest.text = clean_text(rest_text)
                        run_rest.font.name = FONT_BODY
                        run_rest.font.size = Pt(13)
                        run_rest.font.color.rgb = COLOR_GRAY_DARK
                    else:
                        run_val = p_li.add_run()
                        run_val.text = "■ " + clean_text(li.get_text())
                        run_val.font.name = FONT_BODY
                        run_val.font.size = Pt(13)
                        run_val.font.color.rgb = COLOR_GRAY_DARK
            
            # Plain paragraphs (not in lists)
            for p_node in card_node.find_all("p", recursive=False):
                p_p = tf_card.add_paragraph()
                p_p.text = clean_text(p_node.get_text())
                p_p.font.name = FONT_BODY
                p_p.font.size = Pt(13.5)
                p_p.font.color.rgb = COLOR_GRAY_DARK
                p_p.space_after = Pt(8)
                
            # Code snippet (<pre><code>)
            pre_node = card_node.find("pre")
            if pre_node:
                # We can draw a specialized code box inside the card
                code_text = pre_node.get_text()
                # Remove extra blank lines
                code_lines = [line for line in code_text.split('\n') if line.strip() or len(line) < 40]
                code_text_clean = "\n".join(code_lines).strip()
                
                # Draw small code box inside the card
                # We place it in the lower part of the card
                code_box_left = x_pos + Inches(0.25)
                code_box_top = top_y + Inches(1.2)
                code_box_width = col_width - Inches(0.5)
                code_box_height = col_height - Inches(1.5)
                
                code_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, code_box_left, code_box_top, code_box_width, code_box_height)
                code_bg.fill.solid()
                code_bg.fill.fore_color.rgb = COLOR_TECH_BG
                code_bg.line.color.rgb = COLOR_BLACK
                code_bg.line.width = Pt(1.5)
                
                tf_code = code_bg.text_frame
                format_text_frame(tf_code)
                p_code = tf_code.paragraphs[0]
                p_code.text = code_text_clean
                p_code.font.name = FONT_CODE
                p_code.font.size = Pt(11)
                p_code.font.color.rgb = COLOR_BLACK
                
            # Technical explanation box
            tech_box = card_node.find("div", class_="tech-explain-box")
            if tech_box:
                # Let's draw it in the lower area of the card
                tech_left = x_pos + Inches(0.25)
                tech_top = top_y + Inches(3.2)
                tech_width = col_width - Inches(0.5)
                tech_height = Inches(1.5)
                
                t_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tech_left, tech_top, tech_width, tech_height)
                t_bg.fill.solid()
                t_bg.fill.fore_color.rgb = COLOR_TECH_BG
                t_bg.line.color.rgb = COLOR_BLACK
                t_bg.line.width = Pt(1)
                t_bg.line.dash_style = 2 # Dashed line
                
                tf_tech = t_bg.text_frame
                format_text_frame(tf_tech)
                
                strong_t = tech_box.find("strong")
                p_title = tf_tech.paragraphs[0]
                p_title.text = clean_text(strong_t.get_text() if strong_t else "Nota").upper()
                p_title.font.name = FONT_TITLE
                p_title.font.size = Pt(10)
                p_title.font.color.rgb = COLOR_BLACK
                
                p_desc = tech_box.find("p")
                p_body = tf_tech.add_paragraph()
                p_body.text = clean_text(p_desc.get_text() if p_desc else tech_box.get_text())
                p_body.font.name = FONT_BODY
                p_body.font.size = Pt(11)
                p_body.font.color.rgb = COLOR_GRAY_DARK
                
            # Badges row
            badge_row = card_node.find("div", class_="vis-badge-row")
            if badge_row:
                badges = badge_row.find_all("span", class_="vis-badge")
                badge_y = top_y + Inches(1.2)
                badge_x = x_pos + Inches(0.25)
                
                # Draw each badge as a little bordered box
                for b_idx, badge in enumerate(badges):
                    b_text = clean_text(badge.get_text())
                    # Grid pattern for badges
                    bx = badge_x + (b_idx % 2) * Inches(2.6)
                    by = badge_y + (b_idx // 2) * Inches(0.8)
                    
                    b_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(2.4), Inches(0.6))
                    b_card.fill.solid()
                    b_card.fill.fore_color.rgb = COLOR_WHITE
                    b_card.line.color.rgb = COLOR_BLACK
                    b_card.line.width = Pt(2)
                    
                    tf_b = b_card.text_frame
                    format_text_frame(tf_b)
                    tf_b.margin_top = Inches(0.1)
                    tf_b.margin_bottom = Inches(0.1)
                    
                    p_b = tf_b.paragraphs[0]
                    p_b.alignment = PP_ALIGN.CENTER
                    p_b.text = b_text
                    p_b.font.name = FONT_TITLE
                    p_b.font.size = Pt(11)
                    p_b.font.color.rgb = COLOR_BLACK
                    
        # Footer
        footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.6), Inches(12.133), Inches(0.04))
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = COLOR_BLACK
        footer_line.line.fill.background()
        
        counter_box = slide.shapes.add_textbox(Inches(11.233), Inches(6.7), Inches(1.5), Inches(0.4))
        tf_count = counter_box.text_frame
        p_count = tf_count.paragraphs[0]
        p_count.alignment = PP_ALIGN.RIGHT
        p_count.text = f"{index + 1} / {len(slides_html)}"
        p_count.font.name = FONT_TITLE
        p_count.font.size = Pt(12)
        p_count.font.color.rgb = COLOR_BLACK

# Save presentation
output_path = "presentacion_proyecto.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
