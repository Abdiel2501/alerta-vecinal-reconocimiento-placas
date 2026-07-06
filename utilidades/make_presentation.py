from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, sys

# ─── PALETA: FONDO BLANCO, AZUL Y ROJO ──────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_NAVY  = RGBColor(0x0A, 0x1F, 0x5C)   # Azul oscuro (headers, fondos de sección)
BRAND_BLUE = RGBColor(0x1A, 0x5E, 0xC8)   # Azul primario (acentos, botones)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xFF)   # Azul muy claro (fondos de tarjetas)
BRAND_RED  = RGBColor(0xCC, 0x0F, 0x1A)   # Rojo logo
LIGHT_RED  = RGBColor(0xFF, 0xE5, 0xE5)   # Rojo muy claro (fondos de tarjetas)
TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x2E)   # Texto principal
GREY       = RGBColor(0x6B, 0x7A, 0x99)   # Texto secundario
LINE_GREY  = RGBColor(0xD0, 0xD8, 0xE8)   # Líneas divisoras suaves
YELLOW_ACC = RGBColor(0xFF, 0xA5, 0x00)   # Naranja/Amarillo para hitos

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]

base = os.path.dirname(os.path.abspath(__file__))
IMG_LOGO  = os.path.join(base, "logo_project.png")
IMG_PHOTO = os.path.join(base, "realistic_plate_scan.png")

# ─── HELPERS ─────────────────────────────────────────────────────────
def fill_slide(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, line_color=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=14, bold=False,
        color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb

def img(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)

def header_band(slide, title):
    """Banda azul oscura de encabezado + logo pequeño"""
    rect(slide, 0, 0, SLIDE_W, Inches(1.05), DARK_NAVY)
    rect(slide, 0, Inches(1.05), SLIDE_W, Inches(0.06), BRAND_RED)
    txt(slide, title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.72),
        size=26, bold=True, color=WHITE)
    img(slide, IMG_LOGO, Inches(12.0), Inches(0.05), Inches(0.95), Inches(0.95))

# ═════════════════════════════════════════════════════════════════════
# SLIDE 1 – PORTADA
# ═════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
fill_slide(s1, WHITE)

# Franja azul izquierda
rect(s1, 0, 0, Inches(4.6), SLIDE_H, DARK_NAVY)
rect(s1, Inches(4.6), 0, Inches(0.12), SLIDE_H, BRAND_RED)

# Foto realista del carro escaneado en el lado derecho
img(s1, IMG_PHOTO, Inches(4.75), 0, Inches(8.58), SLIDE_H)
# Overlay suave para que el texto de abajo se lea
rect(s1, Inches(4.75), Inches(5.5), Inches(8.58), Inches(2.0), RGBColor(0xFF, 0xFF, 0xFF))

# Logo
img(s1, IMG_LOGO, Inches(0.3), Inches(0.3), Inches(1.3), Inches(1.3))

# Título
txt(s1, "ALERTA VECINAL", Inches(0.25), Inches(1.8), Inches(4.1), Inches(0.65),
    size=18, bold=True, color=BRAND_RED)
txt(s1, "Sistema Inteligente\nde Reconocimiento\nde Matrículas",
    Inches(0.25), Inches(2.35), Inches(4.1), Inches(2.3),
    size=34, bold=True, color=WHITE)
rect(s1, Inches(0.25), Inches(4.65), Inches(3.5), Inches(0.07), BRAND_RED)
txt(s1, "Reporte de Proyecto — Unidad 2",
    Inches(0.25), Inches(4.8), Inches(4.1), Inches(0.5),
    size=14, color=LIGHT_BLUE)
txt(s1, "Jorge Gabriel Heredia Lara\nAbdiel Gerardo Alonso Herrera\nGeovani Coronado Cruz",
    Inches(0.25), Inches(5.5), Inches(4.1), Inches(1.0), size=11, color=LIGHT_BLUE)
txt(s1, "Junio 2026", Inches(0.25), Inches(6.7), Inches(3), Inches(0.4),
    size=11, italic=True, color=RGBColor(0xAA, 0xBB, 0xDD))

# ═════════════════════════════════════════════════════════════════════
# SLIDE 2 – EL PROBLEMA
# ═════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
fill_slide(s2, WHITE)
header_band(s2, "El Problema")

# Stat grande lado izquierdo
rect(s2, Inches(0.4), Inches(1.35), Inches(5.4), Inches(5.55), LIGHT_RED)
rect(s2, Inches(0.4), Inches(1.35), Inches(0.12), Inches(5.55), BRAND_RED)
txt(s2, "63,726", Inches(0.65), Inches(1.55), Inches(5.0), Inches(2.2),
    size=80, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)
txt(s2, "vehículos robados en México\ndurante 2024",
    Inches(0.65), Inches(3.65), Inches(5.0), Inches(0.85),
    size=16, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
txt(s2, "Fuente: AMIS, 2024",
    Inches(0.65), Inches(4.55), Inches(5.0), Inches(0.4),
    size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)
txt(s2, "Solo el 40% de los vehículos\nrobados son recuperados.",
    Inches(0.65), Inches(5.1), Inches(5.0), Inches(0.9),
    size=13, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# Línea divisora
rect(s2, Inches(6.1), Inches(1.5), Inches(0.05), Inches(5.2), LINE_GREY)

# Tarjetas de puntos a la derecha
problems = [
    ("Las camaras solo graban el delito,\nnadie las monitorea en tiempo real."),
    ("Los reportes de robo tardan horas\nen procesarse y llegar a la policia."),
    ("No existe una forma rapida de\nverificar si un vehiculo es robado."),
]
for i, txt_p in enumerate(problems):
    y = Inches(1.5 + i * 1.8)
    rect(s2, Inches(6.4), y, Inches(6.5), Inches(1.5), LIGHT_BLUE)
    rect(s2, Inches(6.4), y, Inches(0.1), Inches(1.5), BRAND_BLUE)
    rect(s2, Inches(6.6), y + Inches(0.35), Inches(0.35), Inches(0.35), BRAND_RED)
    txt(s2, txt_p, Inches(7.1), y + Inches(0.2), Inches(5.6), Inches(1.1),
        size=14, color=TEXT_DARK)

# ═════════════════════════════════════════════════════════════════════
# SLIDE 3 – NUESTRA SOLUCIÓN
# ═════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
fill_slide(s3, WHITE)
header_band(s3, "Nuestra Solucion: Alerta Vecinal")

# Foto realista del carro a la derecha
img(s3, IMG_PHOTO, Inches(7.5), Inches(1.25), Inches(5.6), Inches(3.5))
txt(s3, "Deteccion de placa en tiempo real mediante camara IP con protocolo ONVIF",
    Inches(7.5), Inches(4.8), Inches(5.6), Inches(0.6),
    size=10, italic=True, color=GREY)

# 4 tarjetas de solucion
sol = [
    ("YOLOv8", "Detecta la placa del vehiculo\nen cada fotograma del video.", BRAND_BLUE),
    ("PaddleOCR", "Lee y digitaliza los numeros y\nletras de la placa al instante.", DARK_NAVY),
    ("Flutter App", "Panel de control visual para\noperadores de seguridad.", BRAND_BLUE),
    ("Telegram Bot", "Alerta en el celular en\nmenos de 3 segundos.", BRAND_RED),
]
for i, (title, body, color) in enumerate(sol):
    col = i % 2
    row = i // 2
    x = Inches(0.3 + col * 3.55)
    y = Inches(1.4 + row * 2.85)
    rect(s3, x, y, Inches(3.3), Inches(2.5), LIGHT_BLUE, LINE_GREY)
    rect(s3, x, y, Inches(3.3), Inches(0.12), color)
    txt(s3, title, x + Inches(0.15), y + Inches(0.2), Inches(3.0), Inches(0.55),
        size=17, bold=True, color=color)
    txt(s3, body, x + Inches(0.15), y + Inches(0.8), Inches(3.0), Inches(1.1),
        size=13, color=TEXT_DARK)

# ═════════════════════════════════════════════════════════════════════
# SLIDE 4 – METODOLOGÍA
# ═════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
fill_slide(s4, WHITE)
header_band(s4, "Metodologia: Enfoque Hibrido")

# Gantt lado izquierdo
rect(s4, Inches(0.4), Inches(1.35), Inches(6.0), Inches(5.6), LIGHT_BLUE)
rect(s4, Inches(0.4), Inches(1.35), Inches(6.0), Inches(0.1), BRAND_BLUE)
txt(s4, "Diagrama de GANTT",
    Inches(0.55), Inches(1.5), Inches(5.7), Inches(0.55), size=17, bold=True, color=BRAND_BLUE)
txt(s4, "Para las etapas que deben seguir un orden\nfijo antes de poder continuar con el sistema.",
    Inches(0.55), Inches(2.05), Inches(5.7), Inches(0.7), size=12, italic=True, color=GREY)
gantt_items = [
    ("Recoleccion y etiquetado de datos", "5 dias"),
    ("Entrenamiento del modelo YOLOv8", "7 dias"),
    ("Validacion y ajuste de pesos", "5 dias"),
    ("Integracion del servidor backend", "10 dias"),
]
for i, (task, dur) in enumerate(gantt_items):
    y = Inches(2.95 + i * 0.9)
    rect(s4, Inches(0.55), y + Inches(0.15), Inches(0.25), Inches(0.25), BRAND_BLUE)
    txt(s4, task, Inches(0.95), y, Inches(3.8), Inches(0.55), size=13, color=TEXT_DARK)
    txt(s4, dur, Inches(4.8), y, Inches(1.3), Inches(0.55), size=13, bold=True,
        color=BRAND_BLUE, align=PP_ALIGN.RIGHT)

# Kanban lado derecho
rect(s4, Inches(6.9), Inches(1.35), Inches(6.0), Inches(5.6), LIGHT_RED)
rect(s4, Inches(6.9), Inches(1.35), Inches(6.0), Inches(0.1), BRAND_RED)
txt(s4, "Tablero KANBAN",
    Inches(7.05), Inches(1.5), Inches(5.7), Inches(0.55), size=17, bold=True, color=BRAND_RED)
txt(s4, "Para el desarrollo agil de la App y los modulos\nde mensajeria, con cambios rapidos en equipo.",
    Inches(7.05), Inches(2.05), Inches(5.7), Inches(0.7), size=12, italic=True, color=GREY)

kan_cols = [
    ("Por Hacer",  GREY,       ["Servidor cloud", "Manual de usuario"]),
    ("En Proceso", YELLOW_ACC, ["Calibrar YOLOv8", "App Flutter", "OCR Texto"]),
    ("En Pruebas", BRAND_RED,  ["Latencia Telegram", "Cifrado AES"]),
    ("Listo",      RGBColor(0x00, 0xAA, 0x55), ["BD SQLite", "Graficas App", "Entorno Python"]),
]
for i, (col_title, col_color, items) in enumerate(kan_cols):
    x = Inches(7.05 + i * 1.45)
    rect(s4, x, Inches(2.85), Inches(1.35), Inches(0.38), col_color)
    txt(s4, col_title, x + Inches(0.04), Inches(2.88), Inches(1.27), Inches(0.34),
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items[:3]):
        iy = Inches(3.35 + j * 0.95)
        rect(s4, x + Inches(0.05), iy, Inches(1.25), Inches(0.82),
             WHITE, LINE_GREY)
        txt(s4, item, x + Inches(0.08), iy + Inches(0.08), Inches(1.15), Inches(0.7),
            size=9, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════
# SLIDE 5 – MoSCoW
# ═════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
fill_slide(s5, WHITE)
header_band(s5, "Priorizacion MoSCoW")

moscow = [
    ("MUST HAVE",   BRAND_RED,  LIGHT_RED,  [
        "Deteccion de placas con YOLOv8",
        "Base de datos local SQLite",
        "Alertas automaticas por Telegram",
        "Consulta a BD del REPUVE",
    ]),
    ("SHOULD HAVE", BRAND_BLUE, LIGHT_BLUE, [
        "App visual en Flutter",
        "Alertas por WhatsApp",
        "Scripts de configuracion automatica",
    ]),
    ("COULD HAVE",  RGBColor(0x00, 0x88, 0xCC), RGBColor(0xD6, 0xF0, 0xFF), [
        "Superresolucion ESRGAN para\nplacas borrosas",
        "Bot bidireccional de Telegram",
    ]),
    ("WON'T HAVE",  GREY, RGBColor(0xF0, 0xF0, 0xF4), [
        "Base de datos cloud masiva",
        "Apertura automatica de barreras",
    ]),
]
for i, (label, hdr_color, bg_color, items) in enumerate(moscow):
    x = Inches(0.3 + i * 3.18)
    rect(s5, x, Inches(1.3), Inches(3.0), Inches(5.7), bg_color, LINE_GREY)
    rect(s5, x, Inches(1.3), Inches(3.0), Inches(0.55), hdr_color)
    txt(s5, label, x + Inches(0.1), Inches(1.35), Inches(2.8), Inches(0.48),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        y = Inches(2.05 + j * 1.15)
        rect(s5, x + Inches(0.18), y + Inches(0.18), Inches(0.16), Inches(0.16), hdr_color)
        txt(s5, item, x + Inches(0.45), y, Inches(2.45), Inches(1.1), size=12, color=TEXT_DARK)

# ═════════════════════════════════════════════════════════════════════
# SLIDE 6 – AVANCES AL DÍA DE HOY
# ═════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
fill_slide(s6, WHITE)
header_band(s6, "Avances al Dia de Hoy")

statuses = [
    ("Modelo IA\nEntrenado",        "COMPLETADO", BRAND_BLUE,  LIGHT_BLUE,  DARK_NAVY),
    ("Base de Datos\nSQLite",        "COMPLETADO", BRAND_BLUE,  LIGHT_BLUE,  DARK_NAVY),
    ("Alertas\nTelegram",            "COMPLETADO", BRAND_BLUE,  LIGHT_BLUE,  DARK_NAVY),
    ("App Flutter\nIntegracion",     "EN PROCESO", BRAND_RED,   LIGHT_RED,   BRAND_RED),
]
for i, (label, status, accent, bg, txt_c) in enumerate(statuses):
    x = Inches(0.3 + i * 3.22)
    rect(s6, x, Inches(1.35), Inches(3.0), Inches(4.0), bg, LINE_GREY)
    rect(s6, x, Inches(1.35), Inches(3.0), Inches(0.1), accent)
    # Badge de estado
    rect(s6, x + Inches(0.5), Inches(1.6), Inches(2.0), Inches(0.38), accent)
    txt(s6, status, x + Inches(0.5), Inches(1.63), Inches(2.0), Inches(0.34),
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s6, label, x + Inches(0.1), Inches(2.2), Inches(2.8), Inches(1.3),
        size=20, bold=True, color=txt_c, align=PP_ALIGN.CENTER)

# Barra de progreso
txt(s6, "Progreso general del proyecto:", Inches(0.4), Inches(5.85), Inches(6), Inches(0.45),
    size=14, bold=True, color=TEXT_DARK)
txt(s6, "75%", Inches(11.8), Inches(5.85), Inches(1.1), Inches(0.45),
    size=14, bold=True, color=BRAND_BLUE, align=PP_ALIGN.RIGHT)
rect(s6, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.45), LINE_GREY)
rect(s6, Inches(0.4), Inches(6.4), Inches(9.37), Inches(0.45), BRAND_BLUE)

# Foto pequeña a la derecha
img(s6, IMG_PHOTO, Inches(7.7), Inches(1.45), Inches(5.2), Inches(4.2))

# ═════════════════════════════════════════════════════════════════════
# SLIDE 7 – CRONOGRAMA GANTT
# ═════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
fill_slide(s7, WHITE)
header_band(s7, "Cronograma de Actividades")

gantt_rows = [
    ("Levantamiento de requerimientos",  5,  0,  BRAND_BLUE),
    ("Entrenamiento YOLOv8",             7,  5,  DARK_NAVY),
    ("Desarrollo BD y API",              5,  12, RGBColor(0x00, 0x88, 0x99)),
    ("Programacion App Flutter",         10, 17, BRAND_BLUE),
    ("Integracion Telegram",             3,  27, YELLOW_ACC),
    ("Pruebas y Control de Calidad",     4,  30, BRAND_RED),
    ("Correccion y Entrega Final",       3,  34, RGBColor(0x88, 0x00, 0xCC)),
]
total_days  = 37
chart_left  = Inches(4.4)
chart_width = Inches(8.5)
y_start     = Inches(1.3)
row_h       = Inches(0.78)

# Fondo del área del gráfico
rect(s7, chart_left, y_start, chart_width, row_h * len(gantt_rows), RGBColor(0xF7, 0xF9, 0xFF))

# Líneas de cuadrícula verticales
for d in [0, 5, 10, 15, 20, 25, 30, 35]:
    gx = chart_left + (d / total_days) * chart_width
    rect(s7, gx, y_start, Inches(0.015), row_h * len(gantt_rows), LINE_GREY)
    txt(s7, str(d), gx - Inches(0.15), y_start + row_h * len(gantt_rows) + Inches(0.1),
        Inches(0.4), Inches(0.3), size=9, color=GREY, align=PP_ALIGN.CENTER)

for i, (label, dur, start, color) in enumerate(gantt_rows):
    y = y_start + i * row_h
    # Zebra stripe
    if i % 2 == 0:
        rect(s7, Inches(0.25), y, chart_left - Inches(0.25) + chart_width, row_h,
             RGBColor(0xF0, 0xF4, 0xFF))
    txt(s7, label, Inches(0.3), y + Inches(0.22), Inches(3.9), Inches(0.42), size=11, color=TEXT_DARK)
    bar_x = chart_left + (start / total_days) * chart_width
    bar_w = (dur / total_days) * chart_width
    rect(s7, bar_x, y + Inches(0.18), bar_w, Inches(0.42), color)
    txt(s7, f"{dur}d", bar_x + Inches(0.06), y + Inches(0.2), bar_w, Inches(0.35),
        size=9, bold=True, color=WHITE)

# Hitos
for hito_day, hito_label, hito_color in [(17, "Hito 1\nBackend", BRAND_BLUE),
                                          (30, "Hito 2\nApp OK", YELLOW_ACC),
                                          (37, "Hito 3\nEntrega", BRAND_RED)]:
    hx = chart_left + (hito_day / total_days) * chart_width
    rect(s7, hx - Inches(0.02), y_start, Inches(0.04), row_h * len(gantt_rows), hito_color)
    txt(s7, hito_label, hx - Inches(0.4), y_start + row_h * len(gantt_rows) + Inches(0.4),
        Inches(0.9), Inches(0.55), size=8, bold=True, color=hito_color, align=PP_ALIGN.CENTER)

txt(s7, "Dias habiles — Junio / Julio 2026", chart_left, y_start + row_h * len(gantt_rows) + Inches(0.95),
    chart_width, Inches(0.35), size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════
# SLIDE 8 – RECURSOS E INVERSIÓN
# ═════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
fill_slide(s8, WHITE)
header_band(s8, "Recursos e Inversion")

# Stat grande izq
rect(s8, Inches(0.4), Inches(1.4), Inches(4.5), Inches(5.5), LIGHT_BLUE)
rect(s8, Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.1), BRAND_BLUE)
txt(s8, "Inversion Total", Inches(0.6), Inches(1.6), Inches(4.1), Inches(0.55),
    size=16, bold=True, color=BRAND_BLUE)
txt(s8, "$700", Inches(0.6), Inches(2.25), Inches(4.1), Inches(1.8),
    size=82, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
txt(s8, "MXN", Inches(2.5), Inches(3.9), Inches(2.2), Inches(0.7),
    size=24, bold=True, color=BRAND_RED)
rect(s8, Inches(0.6), Inches(4.6), Inches(4.0), Inches(0.05), LINE_GREY)
txt(s8, "Alta tecnologia, costo accesible.", Inches(0.6), Inches(4.75), Inches(4.1), Inches(0.45),
    size=12, italic=True, color=GREY)
txt(s8, "El resto de las herramientas\nutilizadas son 100% gratuitas\ny de codigo abierto.",
    Inches(0.6), Inches(5.2), Inches(4.1), Inches(0.9), size=12, color=TEXT_DARK)

# Lista de recursos derecha
resources = [
    ("Camara IP protocolo ONVIF",              "$700 MXN",  BRAND_RED),
    ("Python + YOLOv8 + PaddleOCR",           "Gratuito",  BRAND_BLUE),
    ("Flutter SDK + Dart",                     "Gratuito",  BRAND_BLUE),
    ("REPUVE API + SQLite local",              "Publico",   DARK_NAVY),
    ("Bot de Telegram / API mensajeria",       "Gratuito",  BRAND_BLUE),
]
for i, (label, cost, color) in enumerate(resources):
    y = Inches(1.45 + i * 1.07)
    rect(s8, Inches(5.3), y, Inches(7.6), Inches(0.88), LIGHT_BLUE, LINE_GREY)
    rect(s8, Inches(5.3), y, Inches(0.1), Inches(0.88), color)
    txt(s8, label, Inches(5.55), y + Inches(0.17), Inches(5.5), Inches(0.55),
        size=14, color=TEXT_DARK)
    txt(s8, cost, Inches(11.1), y + Inches(0.17), Inches(1.65), Inches(0.55),
        size=14, bold=True, color=color, align=PP_ALIGN.RIGHT)

# ═════════════════════════════════════════════════════════════════════
# SLIDE 9 – PRÓXIMOS PASOS Y CIERRE
# ═════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank)
fill_slide(s9, WHITE)

# Franja azul oscura abajo
rect(s9, 0, Inches(5.8), SLIDE_W, Inches(1.7), DARK_NAVY)
rect(s9, 0, Inches(5.8), SLIDE_W, Inches(0.08), BRAND_RED)

# Foto realista de fondo-derecho
img(s9, IMG_PHOTO, Inches(6.7), Inches(0.8), Inches(6.4), Inches(4.8))
rect(s9, Inches(6.7), Inches(0.8), Inches(6.4), Inches(4.8), RGBColor(0xFF, 0xFF, 0xFF))
img(s9, IMG_PHOTO, Inches(6.7), Inches(0.8), Inches(6.4), Inches(4.8))
# Overlay blanco semitransparente simulado
rect(s9, Inches(6.7), Inches(0.8), Inches(0.08), Inches(4.8), BRAND_BLUE)

# Logo + título
img(s9, IMG_LOGO, Inches(0.3), Inches(0.2), Inches(1.2), Inches(1.2))
txt(s9, "Proximos Pasos", Inches(1.7), Inches(0.35), Inches(5.0), Inches(0.7),
    size=26, bold=True, color=DARK_NAVY)
rect(s9, Inches(1.7), Inches(1.1), Inches(4.7), Inches(0.06), BRAND_RED)

next_steps = [
    ("Calibrar YOLOv8 para condiciones de lluvia y baja luz.",   BRAND_BLUE),
    ("Finalizar integracion Frontend-Backend en App Flutter.",    DARK_NAVY),
    ("Generar reportes exportables a PDF desde la App.",          BRAND_RED),
    ("Pruebas de estres final y ajuste de latencia.",             GREY),
]
for i, (step, color) in enumerate(next_steps):
    y = Inches(1.35 + i * 1.0)
    rect(s9, Inches(1.7), y + Inches(0.2), Inches(0.28), Inches(0.28), color)
    txt(s9, step, Inches(2.15), y, Inches(4.3), Inches(0.85), size=13, color=TEXT_DARK)

# Texto de cierre en franja azul
txt(s9, "github.com/Abdiel2501/yolo-plate-recognition",
    Inches(0.5), Inches(6.1), Inches(9.0), Inches(0.5),
    size=13, color=LIGHT_BLUE, italic=True)
txt(s9, "Abdiel Alonso  |  Jorge Heredia  |  Geovani Coronado  |  Junio 2026",
    Inches(0.5), Inches(6.7), Inches(9.0), Inches(0.5), size=11, color=GREY)
img(s9, IMG_LOGO, Inches(11.8), Inches(5.9), Inches(1.1), Inches(1.1))

# ─── GUARDAR ─────────────────────────────────────────────────────────
out_path = os.path.join(base, "Presentacion_Unidad2_ANPR.pptx")
prs.save(out_path)
print(f"Presentacion guardada en: {out_path}")
